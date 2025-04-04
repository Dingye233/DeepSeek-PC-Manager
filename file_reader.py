import os
import shutil
import tempfile
import mimetypes
from pathlib import Path

# 用于处理不同文件类型的库（这些需要通过pip安装）
try:
    # Office文件处理
    import docx  # 需要安装: pip install python-docx
    import openpyxl  # 需要安装: pip install openpyxl
    from pptx import Presentation  # 需要安装: pip install python-pptx
    
    # PDF文件处理
    from PyPDF2 import PdfReader  # 需要安装: pip install PyPDF2
    
    OFFICE_SUPPORT = True
except ImportError:
    OFFICE_SUPPORT = False
    print("警告: 未安装Office文件支持库，将只能读取基本文本文件。")
    print("可以运行以下命令安装支持: pip install python-docx openpyxl python-pptx PyPDF2")


def read_file(file_path, max_size=1024*1024*10, encoding='utf-8', extract_text_only=True):
    """
    通用文件读取函数，支持多种文件类型
    
    参数:
        file_path (str): 文件路径
        max_size (int): 读取文件的最大大小（字节），默认10MB
        encoding (str): 文本文件的编码，默认utf-8
        extract_text_only (bool): 对于Office文件是否只提取文本内容，默认True
        
    返回:
        dict: 包含文件内容和元数据的字典
    """
    # 检查文件是否存在
    if not os.path.exists(file_path):
        return {
            "success": False,
            "error": f"文件不存在: {file_path}",
            "content": None,
            "file_type": None,
            "size": 0
        }
    
    # 检查文件大小
    file_size = os.path.getsize(file_path)
    if file_size > max_size:
        return {
            "success": False,
            "error": f"文件太大 ({file_size} 字节), 超过了最大限制 ({max_size} 字节)",
            "content": None,
            "file_type": None,
            "size": file_size
        }
    
    # 猜测文件类型
    file_ext = os.path.splitext(file_path)[1].lower()
    content_type, _ = mimetypes.guess_type(file_path)
    
    try:
        # 处理文本文件
        if content_type and content_type.startswith('text/') or file_ext in ['.txt', '.py', '.js', '.html', '.css', '.json', '.xml', '.md', '.csv']:
            return read_text_file(file_path, encoding)
        
        # 处理Office文件和PDF
        elif OFFICE_SUPPORT:
            # Word文档
            if file_ext in ['.docx', '.doc']:
                return read_word_document(file_path, extract_text_only)
            
            # Excel文件
            elif file_ext in ['.xlsx', '.xls']:
                return read_excel_file(file_path, extract_text_only)
            
            # PowerPoint文件
            elif file_ext in ['.pptx', '.ppt']:
                return read_powerpoint(file_path, extract_text_only)
            
            # PDF文件
            elif file_ext == '.pdf':
                return read_pdf_file(file_path)
            
            # 图片文件 - 只能提供元数据，不能提取内容
            elif content_type and content_type.startswith('image/'):
                return {
                    "success": True,
                    "content": f"[图片文件: {os.path.basename(file_path)}]",
                    "file_type": content_type,
                    "size": file_size,
                    "metadata": {
                        "dimensions": get_image_dimensions(file_path)
                    }
                }
            
            # 二进制文件 - 不能提取内容
            else:
                return {
                    "success": True,
                    "content": f"[二进制文件: {os.path.basename(file_path)}]",
                    "file_type": content_type or "application/octet-stream",
                    "size": file_size
                }
        else:
            # 如果没有安装Office支持库，但文件是Office文件
            if file_ext in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.pdf']:
                return {
                    "success": False,
                    "error": f"文件类型 {file_ext} 需要安装额外的库支持",
                    "content": None,
                    "file_type": content_type,
                    "size": file_size
                }
            else:
                # 尝试作为文本文件读取
                return read_text_file(file_path, encoding)
    
    except Exception as e:
        return {
            "success": False,
            "error": f"读取文件时出错: {str(e)}",
            "content": None,
            "file_type": content_type,
            "size": file_size
        }


def read_text_file(file_path, encoding='utf-8'):
    """读取文本文件"""
    try:
        with open(file_path, 'r', encoding=encoding) as file:
            content = file.read()
        
        return {
            "success": True,
            "content": content,
            "file_type": "text/plain",
            "size": os.path.getsize(file_path)
        }
    except UnicodeDecodeError:
        # 如果UTF-8解码失败，尝试其他编码
        encodings = ['gbk', 'gb2312', 'latin-1', 'cp1252']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as file:
                    content = file.read()
                return {
                    "success": True,
                    "content": content,
                    "file_type": "text/plain",
                    "size": os.path.getsize(file_path),
                    "detected_encoding": enc
                }
            except UnicodeDecodeError:
                continue
        
        # 所有尝试都失败，读取为二进制并返回前100个字节的十六进制表示
        with open(file_path, 'rb') as file:
            binary_content = file.read(100)
        
        return {
            "success": False,
            "error": "无法解码文本文件，可能是二进制文件",
            "content": f"二进制数据 (前100字节): {binary_content.hex()}",
            "file_type": "application/octet-stream",
            "size": os.path.getsize(file_path)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取文本文件时出错: {str(e)}",
            "content": None,
            "file_type": "text/plain",
            "size": os.path.getsize(file_path)
        }


def read_word_document(file_path, extract_text_only=True):
    """读取Word文档"""
    try:
        doc = docx.Document(file_path)
        
        if extract_text_only:
            content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
            
            # 添加表格内容
            for table in doc.tables:
                for row in table.rows:
                    content += "\n" + " | ".join([cell.text for cell in row.cells])
        else:
            # 提供更详细的文档结构
            content = {"paragraphs": [], "tables": []}
            
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    content["paragraphs"].append({
                        "text": paragraph.text,
                        "style": paragraph.style.name if paragraph.style else "Normal"
                    })
            
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    table_data.append([cell.text for cell in row.cells])
                content["tables"].append(table_data)
        
        return {
            "success": True,
            "content": content,
            "file_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "size": os.path.getsize(file_path)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取Word文档时出错: {str(e)}",
            "content": None,
            "file_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "size": os.path.getsize(file_path)
        }


def read_excel_file(file_path, extract_text_only=True):
    """读取Excel文件"""
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        
        if extract_text_only:
            content = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content += f"\n\n--- 工作表: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        content += "\n" + " | ".join([str(cell) if cell is not None else "" for cell in row])
        else:
            # 提供结构化的Excel数据
            content = {}
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    sheet_data.append([str(cell) if cell is not None else "" for cell in row])
                
                content[sheet_name] = sheet_data
        
        return {
            "success": True,
            "content": content,
            "file_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "size": os.path.getsize(file_path)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取Excel文件时出错: {str(e)}",
            "content": None, 
            "file_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "size": os.path.getsize(file_path)
        }


def read_powerpoint(file_path, extract_text_only=True):
    """读取PowerPoint文件"""
    try:
        presentation = Presentation(file_path)
        
        if extract_text_only:
            content = ""
            for i, slide in enumerate(presentation.slides):
                content += f"\n\n--- 幻灯片 {i+1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += f"\n{shape.text}"
        else:
            # 提供结构化的PowerPoint数据
            content = []
            for i, slide in enumerate(presentation.slides):
                slide_content = {
                    "slide_number": i+1,
                    "shapes": []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_content["shapes"].append({
                            "type": shape.name,
                            "text": shape.text
                        })
                
                content.append(slide_content)
        
        return {
            "success": True,
            "content": content,
            "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "size": os.path.getsize(file_path)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取PowerPoint文件时出错: {str(e)}",
            "content": None,
            "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "size": os.path.getsize(file_path)
        }


def read_pdf_file(file_path):
    """读取PDF文件"""
    try:
        reader = PdfReader(file_path)
        content = ""
        
        for i, page in enumerate(reader.pages):
            content += f"\n\n--- 页面 {i+1} ---\n"
            try:
                page_text = page.extract_text()
                if page_text:
                    content += page_text
                else:
                    content += "[页面可能包含扫描图像或无法提取的文本]"
            except Exception as e:
                content += f"[无法提取页面文本: {str(e)}]"
        
        return {
            "success": True,
            "content": content,
            "file_type": "application/pdf",
            "size": os.path.getsize(file_path),
            "metadata": {
                "pages": len(reader.pages),
                "title": reader.metadata.title if reader.metadata and hasattr(reader.metadata, 'title') else None,
                "author": reader.metadata.author if reader.metadata and hasattr(reader.metadata, 'author') else None
            }
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取PDF文件时出错: {str(e)}",
            "content": None,
            "file_type": "application/pdf",
            "size": os.path.getsize(file_path)
        }


def get_image_dimensions(file_path):
    """获取图片尺寸（需要安装Pillow库）"""
    try:
        from PIL import Image
        with Image.open(file_path) as img:
            return f"{img.width}x{img.height}"
    except ImportError:
        return "需要安装Pillow库(pip install Pillow)才能获取图片尺寸"
    except Exception:
        return "无法确定"


def copy_file(source_path, destination_path, overwrite=False):
    """
    复制文件到指定位置
    
    参数:
        source_path (str): 源文件路径
        destination_path (str): 目标文件路径
        overwrite (bool): 是否覆盖已存在的文件
        
    返回:
        dict: 操作结果信息
    """
    try:
        # 检查源文件是否存在
        if not os.path.exists(source_path):
            return {
                "success": False,
                "error": f"源文件不存在: {source_path}"
            }
        
        # 检查目标文件是否已存在
        if os.path.exists(destination_path) and not overwrite:
            return {
                "success": False,
                "error": f"目标文件已存在: {destination_path} (设置overwrite=True以覆盖)"
            }
        
        # 确保目标目录存在
        os.makedirs(os.path.dirname(os.path.abspath(destination_path)), exist_ok=True)
        
        # 复制文件
        shutil.copy2(source_path, destination_path)
        
        return {
            "success": True,
            "message": f"文件已复制到: {destination_path}",
            "source": source_path,
            "destination": destination_path
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"复制文件时出错: {str(e)}",
            "source": source_path,
            "destination": destination_path
        }


def write_file(file_path, content, encoding='utf-8', mode='w', make_dirs=True):
    """
    写入内容到文件
    
    参数:
        file_path (str): 文件路径
        content (str): 要写入的内容
        encoding (str): 编码方式(文本模式时使用)
        mode (str): 写入模式('w'表示覆盖，'a'表示追加)
        make_dirs (bool): 如果父目录不存在，是否创建
        
    返回:
        dict: 操作结果
    """
    try:
        # 确保目录存在
        directory = os.path.dirname(os.path.abspath(file_path))
        if make_dirs and not os.path.exists(directory):
            os.makedirs(directory, exist_ok=True)
        
        # 写入内容
        if 'b' in mode:  # 二进制模式
            with open(file_path, mode) as file:
                file.write(content)
        else:  # 文本模式
            with open(file_path, mode, encoding=encoding) as file:
                file.write(content)
        
        return {
            "success": True,
            "message": f"内容已{'追加到' if mode == 'a' else '写入'}: {file_path}",
            "file_path": file_path,
            "size": os.path.getsize(file_path)
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"写入文件时出错: {str(e)}",
            "file_path": file_path
        }


# 示例用法
if __name__ == "__main__":
    # 测试读取文本文件
    text_result = read_file("file_utils.py")
    if text_result["success"]:
        print(f"文本文件读取成功，大小: {text_result['size']}字节")
        print(f"内容前100个字符: {text_result['content'][:100]}...")
    else:
        print(f"读取错误: {text_result['error']}")
    
    # 如果安装了Office支持，可以测试读取Office文件
    if OFFICE_SUPPORT:
        # 假设当前目录有一个测试Word文档
        docx_file = "test.docx"
        if os.path.exists(docx_file):
            docx_result = read_file(docx_file)
            if docx_result["success"]:
                print(f"\nWord文档读取成功，大小: {docx_result['size']}字节")
                print(f"内容前100个字符: {str(docx_result['content'])[:100]}...")
            else:
                print(f"\nWord文档读取错误: {docx_result['error']}") 